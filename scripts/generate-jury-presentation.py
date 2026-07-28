"""Generate the editable SentinelMonitorIA executive jury presentation.

Run from the repository root:
    python scripts/generate-jury-presentation.py

The generated PowerPoint intentionally excludes demo passwords, JWTs and API keys.
All claims are based on evidence already documented in the repository.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "jury" / "SentinelMonitorIA-Presentacion-Jurado-AWS-Codigo-Facilito.pptx"

HERO = ROOT / "Imagenes" / "1769717477540.png"
AWS_LOGO = ROOT / "docs" / "jury" / "logos" / "logo aws.png"
CODIGO_LOGO = ROOT / "Imagenes" / "logo1 codigo facilito.png"
ARCHITECTURE = ROOT / "Diagrama de arquitectura" / "Infraestructura Completo.png"
CAPTURES = ROOT / "docs" / "jury" / "Capturas"

DEMO_URL = "http://sentinelmonitoria-staging-demo-952763303883-20260726060638.s3-website-us-east-1.amazonaws.com"
REPO_URL = "https://github.com/ramonesj/SentinelMonitorIA"
HEALTH_URL = "http://sm-staging-alb-1278334952.us-east-1.elb.amazonaws.com/health"

SLIDE_W = 13.333
SLIDE_H = 7.5
TOTAL_SLIDES = 15
DISPLAY_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Cascadia Mono"

COLORS = {
    "navy": "08111F",
    "navy_2": "0E1A2B",
    "navy_3": "132337",
    "ink": "101C31",
    "teal": "078C73",
    "mint": "79E2C5",
    "mint_soft": "DDF8F0",
    "aws": "FF9900",
    "purple": "7C3AED",
    "blue": "2F80ED",
    "red": "E5485D",
    "amber": "E5A126",
    "green": "16A36A",
    "white": "FFFFFF",
    "paper": "F4F7FB",
    "card": "FFFFFF",
    "line": "D7E1EC",
    "muted": "65758B",
    "muted_dark": "A7B4C6",
    "slate": "26364C",
    "soft_blue": "EAF3FF",
    "soft_purple": "F1EBFF",
    "soft_red": "FFECEF",
    "soft_amber": "FFF5DF",
    "soft_green": "E7F8F1",
}


def C(name_or_hex: str) -> RGBColor:
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor.from_string(value)


def I(value: float):
    return Inches(value)


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C(color)


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str | None = None,
    radius: bool = True,
    line_width: float = 1.0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    if line:
        shape.line.color.rgb = C(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x: float, y: float, w: float, color: str, height: float = 0.025):
    return add_box(slide, x, y, w, height, color, radius=False)


def style_run(run, size: float, color: str, bold: bool = False, font: str = BODY_FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = C(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: str = "ink",
    bold: bool = False,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    line_spacing: float | None = None,
):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = I(margin)
    frame.margin_right = I(margin)
    frame.margin_top = I(margin)
    frame.margin_bottom = I(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    if line_spacing:
        paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    style_run(run, size, color, bold, font)
    return shape


def add_rich_text(
    slide,
    runs: Sequence[tuple[str, float, str, bool, str]],
    x: float,
    y: float,
    w: float,
    h: float,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, size, color, bold, font in runs:
        run = paragraph.add_run()
        run.text = text
        style_run(run, size, color, bold, font)
    return shape


def add_bullets(
    slide,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 15,
    color: str = "ink",
    bullet_color: str = "teal",
    spacing: float = 7,
):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = I(0.03)
    frame.margin_top = frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(spacing)
        paragraph.line_spacing = 1.06
        bullet = paragraph.add_run()
        bullet.text = "●  "
        style_run(bullet, max(8, size - 4), bullet_color, True, BODY_FONT)
        run = paragraph.add_run()
        run.text = item
        style_run(run, size, color, False, BODY_FONT)
    return shape


def add_chip(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    fill: str,
    color: str,
    line: str | None = None,
    size: float = 10.5,
    hyperlink: str | None = None,
):
    shape = add_box(slide, x, y, w, 0.34, fill, line or fill, radius=True, line_width=0.8)
    add_text(slide, text, x, y + 0.005, w, 0.31, size, color, True, BODY_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    if hyperlink:
        shape.click_action.hyperlink.address = hyperlink
    return shape


def add_icon_circle(slide, label: str, x: float, y: float, diameter: float, fill: str, color: str = "white", size: float = 14):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(diameter), I(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    shape.line.fill.background()
    add_text(slide, label, x, y, diameter, diameter, size, color, True, BODY_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return shape


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        image_w, image_h = image.size
    source_ratio = image_w / image_h
    target_ratio = w / h
    if source_ratio >= target_ratio:
        draw_w = w
        draw_h = w / source_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * source_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    return slide.shapes.add_picture(str(path), I(draw_x), I(draw_y), I(draw_w), I(draw_h))


def add_picture_cover(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    centering: tuple[float, float] = (0.5, 0.5),
    pixels_per_inch: int = 180,
):
    target_size = (max(400, int(w * pixels_per_inch)), max(240, int(h * pixels_per_inch)))
    with Image.open(path) as source:
        image = ImageOps.fit(
            source.convert("RGB"),
            target_size,
            method=Image.Resampling.LANCZOS,
            centering=centering,
        )
        stream = BytesIO()
        image.save(stream, format="JPEG", quality=91, optimize=True, progressive=True)
    stream.seek(0)
    return slide.shapes.add_picture(stream, I(x), I(y), I(w), I(h))


def add_logo_pill(slide, path: Path, x: float, y: float, w: float, h: float):
    add_box(slide, x, y, w, h, "white", "line", radius=True, line_width=0.6)
    add_picture_contain(slide, path, x + 0.12, y + 0.08, w - 0.24, h - 0.16)


def add_title(slide, eyebrow: str, title: str, subtitle: str | None = None, dark: bool = False):
    primary = "white" if dark else "ink"
    secondary = "muted_dark" if dark else "muted"
    add_text(slide, eyebrow.upper(), 0.62, 0.38, 7.8, 0.22, 9.5, "mint" if dark else "teal", True, MONO_FONT)
    add_text(slide, title, 0.62, 0.68, 11.9, 0.58, 27, primary, True, DISPLAY_FONT)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.23, 11.7, 0.36, 13.5, secondary, False, BODY_FONT)


def add_footer(slide, slide_number: int, dark: bool = False):
    muted = "muted_dark" if dark else "muted"
    line = "slate" if dark else "line"
    add_line(slide, 0.62, 7.06, 12.08, line, 0.012)
    add_text(slide, "SENTINELMONITORIA · PRESENTACIÓN EJECUTIVA", 0.62, 7.13, 5.8, 0.16, 7.6, muted, True, MONO_FONT)
    add_text(slide, f"{slide_number:02d} / {TOTAL_SLIDES:02d}", 11.95, 7.13, 0.75, 0.16, 7.6, muted, True, MONO_FONT, PP_ALIGN.RIGHT)


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def add_status_dot(slide, x: float, y: float, color: str, diameter: float = 0.10):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(diameter), I(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(color)
    shape.line.fill.background()
    return shape


def add_arrow(slide, x: float, y: float, w: float, h: float = 0.22, color: str = "mint"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(color)
    shape.line.fill.background()
    return shape


def add_metric_card(slide, value: str, label: str, detail: str, x: float, y: float, w: float, fill: str = "white", accent: str = "teal"):
    add_box(slide, x, y, w, 1.12, fill, "line", radius=True)
    add_line(slide, x, y, 0.08, accent, 1.12)
    add_text(slide, value, x + 0.23, y + 0.12, w - 0.35, 0.38, 23, "ink", True, DISPLAY_FONT)
    add_text(slide, label.upper(), x + 0.23, y + 0.55, w - 0.35, 0.18, 8.6, accent, True, MONO_FONT)
    add_text(slide, detail, x + 0.23, y + 0.78, w - 0.35, 0.20, 9.5, "muted")


def add_step_card(slide, number: str, title: str, detail: str, x: float, y: float, w: float, accent: str):
    add_box(slide, x, y, w, 1.10, "navy_2", "slate", radius=True)
    add_icon_circle(slide, number, x + 0.12, y + 0.16, 0.38, accent, "white", 10)
    add_text(slide, title, x + 0.60, y + 0.13, w - 0.72, 0.30, 12, "white", True, DISPLAY_FONT)
    add_text(slide, detail, x + 0.60, y + 0.48, w - 0.72, 0.43, 9.2, "muted_dark")


def add_small_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, accent: str, dark: bool = False):
    fill = "navy_2" if dark else "white"
    line = "slate" if dark else "line"
    title_color = "white" if dark else "ink"
    body_color = "muted_dark" if dark else "muted"
    add_box(slide, x, y, w, h, fill, line, radius=True)
    add_line(slide, x, y, w, accent, 0.055)
    add_text(slide, title, x + 0.20, y + 0.20, w - 0.40, 0.30, 13, title_color, True, DISPLAY_FONT)
    add_text(slide, body, x + 0.20, y + 0.58, w - 0.40, h - 0.72, 10.5, body_color, False, BODY_FONT)


def new_slide(prs: Presentation, color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, color)
    return slide


def build_presentation() -> Presentation:
    required = [HERO, AWS_LOGO, CODIGO_LOGO, ARCHITECTURE]
    required.extend(CAPTURES / f"Screenshot_{number}.png" for number in (4, 5, 10, 14, 15, 17))
    missing = [str(path) for path in required if not path.exists()]
    if missing:
        raise FileNotFoundError("Missing presentation assets:\n" + "\n".join(missing))

    prs = Presentation()
    prs.slide_width = I(SLIDE_W)
    prs.slide_height = I(SLIDE_H)
    properties = prs.core_properties
    properties.title = "SentinelMonitorIA · Presentación ejecutiva para jurados"
    properties.subject = "Observabilidad, AIOps y arquitectura desplegada en AWS"
    properties.author = "Equipo SentinelMonitorIA"
    properties.last_modified_by = "Kiro"
    properties.keywords = "SentinelMonitorIA, AWS, AIOps, observabilidad, Lex V2, ECS, FastAPI"
    properties.comments = "Generada desde evidencias versionadas del repositorio. No contiene secretos de demostración."

    # 01 — Cover
    slide = new_slide(prs, "navy")
    add_picture_cover(slide, HERO, 7.25, 0, 6.083, 7.5, centering=(0.54, 0.5), pixels_per_inch=200)
    add_box(slide, 6.96, 0, 0.30, 7.5, "teal", radius=False)
    add_logo_pill(slide, AWS_LOGO, 0.68, 0.48, 1.32, 0.55)
    add_logo_pill(slide, CODIGO_LOGO, 2.14, 0.48, 2.20, 0.55)
    add_text(slide, "PRESENTACIÓN EJECUTIVA · JURADO 2026", 0.70, 1.50, 5.8, 0.26, 10, "mint", True, MONO_FONT)
    add_text(slide, "SentinelMonitorIA", 0.68, 1.91, 6.25, 0.75, 38, "white", True, DISPLAY_FONT)
    add_text(slide, "De telemetría dispersa a\ndecisiones operativas explicables", 0.70, 2.70, 5.95, 1.16, 24, "white", False, DISPLAY_FONT)
    add_text(slide, "Observabilidad y AIOps multi-tenant desplegados en Amazon Web Services.", 0.72, 4.03, 5.65, 0.62, 14.5, "muted_dark")
    add_chip(slide, "AWS DESPLEGADO", 0.72, 4.89, 1.75, "aws", "navy")
    add_chip(slide, "STAGING OPERATIVO", 2.61, 4.89, 1.90, "teal", "white")
    add_chip(slide, "us-east-1", 4.65, 4.89, 1.15, "navy_3", "mint", "slate")
    add_status_dot(slide, 0.72, 5.76, "mint", 0.11)
    add_text(slide, "Pipeline asíncrono · Alertas explicables · Chat Lex V2 · IaC auditable", 0.94, 5.69, 5.70, 0.38, 11.5, "muted_dark")
    add_text(slide, "AWS · Código Facilito · Kiro", 0.70, 6.74, 4.8, 0.20, 8.5, "muted_dark", True, MONO_FONT)
    add_text(slide, "01 / 15", 6.06, 6.74, 0.60, 0.20, 8.5, "muted_dark", True, MONO_FONT, PP_ALIGN.RIGHT)
    add_notes(slide, """
Abrir con una idea simple: SentinelMonitorIA convierte señales técnicas en decisiones operativas claras. No presentamos únicamente un diagrama; mostramos un MVP desplegado en AWS, con un flujo real desde la ingesta hasta la alerta y una interfaz que el jurado puede recorrer.
""")

    # 02 — Problem
    slide = new_slide(prs, "paper")
    add_title(slide, "01 · El problema", "Las operaciones no fallan por falta de datos.", "Fallan cuando las señales no se convierten a tiempo en contexto y decisiones.")
    add_text(slide, "Métricas, logs y eventos llegan en paralelo, pero el operador necesita una sola historia: qué ocurrió, por qué importa y cuál es el siguiente paso seguro.", 0.68, 1.83, 4.55, 1.55, 21, "ink", True, DISPLAY_FONT)
    add_box(slide, 0.68, 3.62, 4.55, 1.74, "navy", "navy", radius=True)
    add_text(slide, "EL COSTE OPERATIVO", 0.96, 3.93, 2.3, 0.22, 9.5, "mint", True, MONO_FONT)
    add_text(slide, "Más ruido · más contexto manual · respuestas más lentas", 0.96, 4.30, 3.92, 0.65, 18, "white", True, DISPLAY_FONT)
    problems = [
        ("01", "Sobrecarga de señales", "Múltiples fuentes sin una narrativa común ni priorización.", "blue"),
        ("02", "Procesamiento acoplado", "El emisor no debería esperar persistencia y análisis pesado.", "purple"),
        ("03", "Alertas sin explicación", "Una alarma aislada no entrega evidencia ni recomendaciones.", "red"),
        ("04", "Contexto fragmentado", "Identidad, organización y operación deben permanecer aisladas.", "teal"),
    ]
    for idx, (number, title, body, accent) in enumerate(problems):
        col = idx % 2
        row = idx // 2
        x = 5.58 + col * 3.55
        y = 1.82 + row * 2.05
        add_box(slide, x, y, 3.25, 1.73, "white", "line", radius=True)
        add_icon_circle(slide, number, x + 0.22, y + 0.22, 0.44, accent, "white", 10)
        add_text(slide, title, x + 0.82, y + 0.19, 2.15, 0.38, 13.5, "ink", True, DISPLAY_FONT)
        add_text(slide, body, x + 0.22, y + 0.77, 2.80, 0.68, 10.8, "muted")
    add_box(slide, 5.58, 6.08, 6.80, 0.57, "mint_soft", "mint", radius=True)
    add_text(slide, "OBJETIVO  →  reducir incertidumbre sin perder trazabilidad, seguridad ni control de costes", 5.78, 6.20, 6.40, 0.28, 10.6, "teal", True, MONO_FONT, PP_ALIGN.CENTER)
    add_footer(slide, 2)
    add_notes(slide, """
El problema no es recolectar más datos. El problema es correlacionarlos con velocidad y conservar el contexto de organización, evidencia y seguridad. SentinelMonitorIA se enfoca en cerrar esa distancia entre señal y decisión, sin bloquear al productor ni depender de una caja negra generativa.
""")

    # 03 — Solution
    slide = new_slide(prs, "navy")
    add_title(slide, "02 · La solución", "Una historia operativa trazable, de extremo a extremo.", "Cada etapa queda desacoplada, observable y preparada para evolucionar.", dark=True)
    steps = [
        ("1", "Ingerir", "Métricas, logs y eventos"),
        ("2", "Aceptar", "HTTP 202 inmediato"),
        ("3", "Desacoplar", "Redis Streams"),
        ("4", "Persistir", "PostgreSQL privado"),
        ("5", "Analizar", "Reglas explicables"),
        ("6", "Decidir", "Alert + Lex V2"),
    ]
    step_w = 1.78
    for idx, (number, title, detail) in enumerate(steps):
        x = 0.48 + idx * 2.12
        add_step_card(slide, number, title, detail, x, 2.03, step_w, "teal" if idx not in (1, 4) else "aws")
        if idx < len(steps) - 1:
            add_arrow(slide, x + step_w + 0.06, 2.46, 0.27, 0.20, "mint")
    add_box(slide, 0.62, 3.58, 12.08, 0.72, "navy_3", "slate", radius=True)
    add_rich_text(
        slide,
        [
            ("Resultado demostrable: ", 14, "mint", True, BODY_FONT),
            ("un batch real aceptado por la API termina como ", 14, "white", False, BODY_FONT),
            ("AIAnalysis → Alert", 14, "aws", True, MONO_FONT),
            (" y puede consultarse desde el dashboard o Lex V2.", 14, "white", False, BODY_FONT),
        ],
        0.92,
        3.82,
        11.50,
        0.28,
        PP_ALIGN.CENTER,
    )
    pillars = [
        ("ASÍNCRONO", "El productor recibe respuesta sin esperar el procesamiento completo.", "blue"),
        ("EXPLICABLE", "Reglas, hallazgos, recomendaciones y analysis_id persistidos.", "purple"),
        ("AISLADO", "JWT, scopes y consultas filtradas por organización.", "teal"),
        ("REPRODUCIBLE", "23 fases CloudFormation y contratos verificables.", "aws"),
    ]
    for idx, (title, body, accent) in enumerate(pillars):
        add_small_card(slide, title, body, 0.62 + idx * 3.08, 4.75, 2.78, 1.45, accent, dark=True)
    add_footer(slide, 3, dark=True)
    add_notes(slide, """
La propuesta se resume en seis verbos. Ingerimos, respondemos 202, desacoplamos con Streams, persistimos, analizamos y entregamos una decisión. El valor diferencial es la trazabilidad completa: la alerta no está pintada en el frontend; nace de telemetría aceptada y procesada por los workers.
""")

    # 04 — Evidence
    slide = new_slide(prs, "paper")
    add_title(slide, "03 · Evidencia operativa", "No es sólo arquitectura: está desplegado en AWS.", "Captura de validación del staging; los contadores evolucionan con la telemetría.")
    add_metric_card(slide, "202", "Ingesta", "Accepted sin bloqueo", 0.62, 1.74, 1.78, "white", "teal")
    add_metric_card(slide, "3 / 3", "Core services", "PostgreSQL · Redis · telemetry", 0.62, 3.03, 1.78, "white", "green")
    add_metric_card(slide, "279", "Eventos", "Snapshot del dashboard", 2.58, 1.74, 1.78, "white", "blue")
    add_metric_card(slide, "100%", "Success rate", "Snapshot del pipeline", 2.58, 3.03, 1.78, "white", "purple")
    add_box(slide, 0.62, 4.48, 3.74, 1.57, "navy", "navy", radius=True)
    add_text(slide, "EVIDENCIA VERIFICADA", 0.90, 4.75, 2.8, 0.23, 9.5, "mint", True, MONO_FONT)
    add_text(slide, "Health 200 · Swagger 200\nAlertas high · Lex intent 0.9", 0.90, 5.10, 3.05, 0.63, 14.5, "white", True, DISPLAY_FONT)
    add_box(slide, 4.70, 1.70, 7.98, 4.66, "white", "line", radius=True)
    add_picture_cover(slide, CAPTURES / "Screenshot_15.png", 4.82, 1.82, 7.74, 4.42, centering=(0.50, 0.47), pixels_per_inch=210)
    add_chip(slide, "STAGING REAL", 10.86, 1.92, 1.42, "green", "white")
    add_text(slide, "Frontend S3 Website · API ALB · backend y workers ECS/Fargate", 4.82, 6.48, 7.74, 0.25, 10.2, "muted", False, MONO_FONT, PP_ALIGN.CENTER)
    add_footer(slide, 4)
    add_notes(slide, """
Esta diapositiva separa evidencia de promesa. El staging respondió 202 para ingesta y 200 en health, Swagger y métricas. La captura muestra los tres servicios principales saludables y datos procesados. Los números 279 y 100 por ciento son una fotografía documentada, no una promesa de rendimiento permanente.
""")

    # 05 — Architecture
    slide = new_slide(prs, "paper")
    add_title(slide, "04 · Arquitectura AWS", "Núcleo privado, borde público controlado y evolución explícita.", "La imagen diferencia el flujo validado de los componentes pendientes o bloqueados.")
    add_box(slide, 0.42, 1.63, 12.50, 5.17, "white", "line", radius=True)
    add_picture_contain(slide, ARCHITECTURE, 0.54, 1.75, 12.26, 4.82)
    add_chip(slide, "ACTIVO", 9.42, 1.76, 0.86, "soft_green", "green", "green", 8.6)
    add_chip(slide, "DEGRADADO", 10.38, 1.76, 1.18, "soft_amber", "amber", "amber", 8.6)
    add_chip(slide, "OBJETIVO", 11.66, 1.76, 1.00, "soft_blue", "blue", "blue", 8.6)
    add_footer(slide, 5)
    add_notes(slide, """
El frontend temporal vive en S3 Website y la API entra por el ALB. Las tareas ECS, PostgreSQL y Redis permanecen privadas. Redis Streams conecta backend, telemetry worker y AI worker. Lex V2 está activo; Bedrock y notificaciones externas aparecen marcados como no disponibles. El diagrama también muestra la evolución objetivo, pero no la confunde con el staging validado.
""")

    # 06 — Pipeline
    slide = new_slide(prs, "navy")
    add_title(slide, "05 · Pipeline asíncrono", "El 202 desacopla al productor del análisis pesado.", "Redis Streams absorbe el trabajo y permite que cada consumer avance a su ritmo.", dark=True)
    top_steps = [
        ("1", "EC2 producer", "API key · telemetry:write", "aws"),
        ("2", "ALB :80", "Punto de entrada API", "blue"),
        ("3", "FastAPI", "Valida y publica batch", "teal"),
        ("4", "Redis telemetry", "Stream + consumer groups", "purple"),
        ("5", "Telemetry worker", "Persiste y confirma", "green"),
    ]
    for idx, (number, title, detail, accent) in enumerate(top_steps):
        x = 0.43 + idx * 2.55
        add_step_card(slide, number, title, detail, x, 1.90, 2.15, accent)
        if idx < 4:
            add_arrow(slide, x + 2.18, 2.34, 0.27, 0.19, "mint")
    add_box(slide, 0.64, 3.32, 2.30, 0.58, "aws", "aws", radius=True)
    add_text(slide, "HTTP 202 ACCEPTED", 0.64, 3.36, 2.30, 0.47, 12, "navy", True, MONO_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(slide, "La conexión HTTP termina aquí; el procesamiento continúa de forma independiente.", 3.18, 3.39, 8.85, 0.30, 12, "muted_dark", False, BODY_FONT)
    bottom_steps = [
        ("6", "RDS + ai_analysis", "Métricas, logs, eventos", "blue"),
        ("7", "AI worker", "Reglas + fallback seguro", "purple"),
        ("8", "AIAnalysis + Alert", "Hallazgos y recomendaciones", "red"),
        ("9", "Dashboard + Lex", "Consulta org-scoped", "teal"),
    ]
    for idx, (number, title, detail, accent) in enumerate(bottom_steps):
        x = 1.12 + idx * 3.02
        add_step_card(slide, number, title, detail, x, 4.30, 2.58, accent)
        if idx < 3:
            add_arrow(slide, x + 2.63, 4.74, 0.27, 0.19, "mint")
    add_chip(slide, "DEDUP · organización + batch_id", 4.58, 5.72, 3.84, "navy_3", "mint", "slate", 10)
    add_footer(slide, 6, dark=True)
    add_notes(slide, """
El productor se autentica con un scope mínimo. FastAPI publica el batch y devuelve 202. Desde ese momento telemetry worker y AI worker operan fuera de la conexión HTTP. PostgreSQL conserva la evidencia y la clave organización más batch_id evita duplicados. El resultado se consulta desde dashboard o Lex sin acoplar la experiencia al productor.
""")

    # 07 — Explainable AIOps
    slide = new_slide(prs, "paper")
    add_title(slide, "06 · AIOps explicable", "Detección útil incluso sin un modelo generativo.", "Reglas determinísticas crean evidencia, recomendaciones y alertas auditables.")
    add_small_card(slide, "CPU / MEMORIA", "Dispara cuando la unidad percent alcanza o supera 90%.", 0.62, 1.76, 3.35, 1.12, "red")
    add_small_card(slide, "LOGS / EVENTOS", "Log error/fatal o evento high/critical.", 0.62, 3.02, 3.35, 1.12, "purple")
    add_small_card(slide, "SALIDA EXPLICABLE", "AIAnalysis, Alert, rule_id, analysis_id, hallazgos y próximos pasos.", 0.62, 4.28, 3.35, 1.32, "teal")
    add_box(slide, 0.62, 5.78, 3.35, 0.58, "soft_amber", "amber", radius=True)
    add_text(slide, "AI_ENABLE_ACTIONS = false", 0.62, 5.84, 3.35, 0.42, 10.5, "amber", True, MONO_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_box(slide, 4.26, 1.73, 8.42, 4.86, "white", "line", radius=True)
    add_picture_cover(slide, CAPTURES / "Screenshot_10.png", 4.38, 1.85, 8.18, 4.62, centering=(0.48, 0.43), pixels_per_inch=220)
    add_chip(slide, "CPU 96%", 4.62, 5.91, 1.16, "soft_red", "red", "red", 8.8)
    add_chip(slide, "MEM 94%", 5.90, 5.91, 1.16, "soft_amber", "amber", "amber", 8.8)
    add_chip(slide, "LOG ERROR", 7.18, 5.91, 1.30, "soft_purple", "purple", "purple", 8.8)
    add_chip(slide, "ALERTA HIGH", 8.60, 5.91, 1.44, "soft_red", "red", "red", 8.8)
    add_footer(slide, 7)
    add_notes(slide, """
La inteligencia actual no depende de Bedrock. El detector evalúa umbrales y severidades verificables. En este incidente controlado aparecen CPU al 96 por ciento, memoria al 94 y un log de error. El sistema conserva la regla, el análisis y recomendaciones de sólo lectura. Las acciones automáticas permanecen deshabilitadas.
""")

    # 08 — Lex
    slide = new_slide(prs, "navy")
    add_title(slide, "07 · Conversación operativa", "Lex V2 entiende la intención; el backend controla la respuesta.", "Español latinoamericano, lectura segura y aislamiento por organización.", dark=True)
    add_text(slide, "AMAZON LEX V2", 0.70, 1.80, 3.5, 0.24, 10, "mint", True, MONO_FONT)
    add_text(slide, "Preguntar por alertas, criticidad o salud sin abandonar el contexto operativo.", 0.70, 2.18, 5.05, 0.92, 21, "white", True, DISPLAY_FONT)
    add_bullets(
        slide,
        [
            "Locale es_419 y cinco intenciones operativas.",
            "OpenAlertsIntent observado con confianza 0.9.",
            "JWT y organización validados en el backend.",
            "Respuesta estructurada, auditable y de sólo lectura.",
        ],
        0.70,
        3.37,
        5.10,
        1.72,
        12.2,
        "muted_dark",
        "mint",
        5,
    )
    add_box(slide, 0.70, 5.35, 5.10, 0.72, "navy_3", "purple", radius=True)
    add_rich_text(
        slide,
        [
            ("Lex interpreta intención", 12, "mint", True, BODY_FONT),
            ("  ≠  ", 12, "muted_dark", True, BODY_FONT),
            ("Bedrock genera contenido", 12, "aws", True, BODY_FONT),
        ],
        0.92,
        5.59,
        4.66,
        0.26,
        PP_ALIGN.CENTER,
    )
    add_box(slide, 6.20, 1.72, 2.88, 4.82, "white", "slate", radius=True)
    add_picture_cover(slide, CAPTURES / "Screenshot_4.png", 6.31, 1.84, 2.66, 4.58, centering=(0.5, 0.49), pixels_per_inch=220)
    add_box(slide, 9.46, 1.72, 2.88, 4.82, "white", "slate", radius=True)
    add_picture_cover(slide, CAPTURES / "Screenshot_5.png", 9.57, 1.84, 2.66, 4.58, centering=(0.5, 0.49), pixels_per_inch=220)
    add_footer(slide, 8, dark=True)
    add_notes(slide, """
Lex V2 reconoce la intención en español latinoamericano. El backend sigue siendo responsable de JWT, permisos, organización y acceso a datos. Es importante no confundir Lex con Bedrock: el chat estructurado está validado aunque el modelo generativo no esté autorizado. La interfaz comunica además que las acciones automáticas están deshabilitadas.
""")

    # 09 — Security
    slide = new_slide(prs, "paper")
    add_title(slide, "08 · Seguridad y multi-tenancy", "Capas de control, mínimo privilegio y límites visibles.", "El staging conserva el núcleo privado y documenta honestamente la concesión HTTP temporal.")
    security_cards = [
        ("IDENTIDAD", "JWT para operadores\nAPI key telemetry:write\nIAM Task Roles", "blue", "ID"),
        ("AISLAMIENTO", "Consultas org-scoped\nECS sin IP pública\nRDS y Redis privados", "teal", "ORG"),
        ("SECRETOS", "Secrets Manager\nTLS para Redis\nSin access keys en imágenes", "purple", "SEC"),
        ("OPERACIÓN", "SSM sin SSH\nProductor no root\nLímites CPU y memoria", "aws", "OPS"),
    ]
    for idx, (title, body, accent, icon) in enumerate(security_cards):
        x = 0.62 + idx * 3.08
        add_box(slide, x, 1.82, 2.78, 2.25, "white", "line", radius=True)
        add_icon_circle(slide, icon, x + 0.20, 2.05, 0.58, accent, "white", 9)
        add_text(slide, title, x + 0.92, 2.08, 1.62, 0.30, 12.5, "ink", True, DISPLAY_FONT)
        add_text(slide, body, x + 0.20, 2.76, 2.34, 0.94, 10.6, "muted")
    add_box(slide, 0.62, 4.44, 12.08, 1.16, "navy", "navy", radius=True)
    network = [
        ("Internet", "white", "navy_3"),
        ("S3 / ALB", "navy", "aws"),
        ("ECS privado", "white", "teal"),
        ("RDS + Redis TLS", "white", "purple"),
    ]
    for idx, (label, text_color, fill) in enumerate(network):
        x = 0.94 + idx * 2.92
        add_box(slide, x, 4.78, 2.18, 0.50, fill, fill, radius=True)
        add_text(slide, label, x, 4.82, 2.18, 0.40, 10.5, text_color, True, MONO_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        if idx < 3:
            add_arrow(slide, x + 2.30, 4.93, 0.35, 0.18, "mint")
    add_box(slide, 0.62, 5.86, 7.72, 0.60, "mint_soft", "mint", radius=True)
    add_text(slide, "AI_ENABLE_ACTIONS=false  ·  notification worker 0/0  ·  sin efectos externos", 0.80, 5.96, 7.36, 0.35, 10.4, "teal", True, MONO_FONT, PP_ALIGN.CENTER)
    add_box(slide, 8.57, 5.86, 4.13, 0.60, "soft_amber", "amber", radius=True)
    add_text(slide, "STAGING: S3 + ALB usan HTTP temporal", 8.70, 5.96, 3.87, 0.35, 9.8, "amber", True, MONO_FONT, PP_ALIGN.CENTER)
    add_footer(slide, 9)
    add_notes(slide, """
La seguridad se explica por capas. Los operadores usan JWT; el productor sólo posee telemetry:write. Los datos se filtran por organización. ECS, RDS y Redis están en subredes privadas, Redis usa TLS y Secrets Manager entrega credenciales en runtime. SSM reemplaza SSH. También declaramos el límite: el borde del staging usa HTTP y debe evolucionar a HTTPS.
""")

    # 10 — Platform / IaC
    slide = new_slide(prs, "navy")
    add_title(slide, "09 · Plataforma cloud", "Servicios activos, degradados y futuros sin ambigüedad.", "La infraestructura se despliega en 23 fases CloudFormation numeradas 00–22.", dark=True)
    add_box(slide, 0.62, 1.72, 7.36, 4.83, "navy_2", "slate", radius=True)
    add_text(slide, "ESTADO DEL STAGING", 0.92, 1.98, 3.0, 0.22, 9.5, "mint", True, MONO_FONT)
    rows = [
        ("Frontend + API", "S3 Website · ALB HTTP", "ACTIVO", "green"),
        ("Cómputo", "backend · telemetry · AI  1/1", "ACTIVO", "green"),
        ("Datos", "RDS PostgreSQL · Redis TLS", "ACTIVO", "green"),
        ("Conversación", "Amazon Lex V2 · es_419", "ACTIVO", "green"),
        ("Bedrock / RAG", "Nova Lite · Titan embeddings", "BLOQUEADO", "amber"),
        ("Notificaciones", "worker desired/running 0/0", "APAGADO", "muted"),
        ("Borde producción", "CloudFront · HTTPS · WAF · HA", "OBJETIVO", "blue"),
    ]
    for idx, (layer, detail, status, status_color) in enumerate(rows):
        y = 2.40 + idx * 0.52
        if idx:
            add_line(slide, 0.92, y - 0.10, 6.76, "slate", 0.01)
        add_text(slide, layer, 0.94, y, 1.72, 0.24, 10.5, "white", True, BODY_FONT)
        add_text(slide, detail, 2.68, y, 3.52, 0.24, 9.7, "muted_dark", False, BODY_FONT)
        add_status_dot(slide, 6.34, y + 0.06, status_color, 0.10)
        add_text(slide, status, 6.52, y, 1.08, 0.24, 8.5, status_color, True, MONO_FONT, PP_ALIGN.RIGHT)
    add_box(slide, 8.28, 1.72, 4.42, 4.83, "navy_2", "slate", radius=True)
    add_text(slide, "CLOUDFORMATION 00–22", 8.58, 1.98, 3.55, 0.22, 9.5, "aws", True, MONO_FONT)
    phase_groups = [
        ("00–02", "Red · NAT · Security Groups", "blue"),
        ("03–04", "IAM · ECR", "teal"),
        ("05–08", "RDS · Redis · Secrets · Logs", "purple"),
        ("09–13", "ALB · ECS · Frontend", "green"),
        ("14–18", "CloudFront · DNS · TLS", "amber"),
        ("19–22", "RAG · AI · Notificaciones", "aws"),
    ]
    for idx, (number, label, accent) in enumerate(phase_groups):
        y = 2.46 + idx * 0.59
        add_box(slide, 8.58, y, 0.86, 0.38, accent, accent, radius=True)
        add_text(slide, number, 8.58, y + 0.01, 0.86, 0.34, 9.3, "white" if accent != "aws" else "navy", True, MONO_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(slide, label, 9.64, y + 0.02, 2.58, 0.34, 10.1, "muted_dark", False, BODY_FONT, valign=MSO_ANCHOR.MIDDLE)
    add_box(slide, 8.58, 6.05, 3.82, 0.30, "navy_3", "slate", radius=True)
    add_text(slide, "EVOLUCIONAR SIN RECONSTRUIR", 8.58, 6.07, 3.82, 0.24, 8.5, "mint", True, MONO_FONT, PP_ALIGN.CENTER)
    add_footer(slide, 10, dark=True)
    add_notes(slide, """
Esta tabla evita sobreprometer. El núcleo de frontend, API, cómputo, datos y Lex está activo. Bedrock está bloqueado y notificaciones apagadas. CloudFront, HTTPS, WAF y alta disponibilidad pertenecen al objetivo de producción. CloudFormation divide el avance en 23 fases auditables, por lo que cada capacidad puede promoverse de forma controlada.
""")

    # 11 — Quality
    slide = new_slide(prs, "paper")
    add_title(slide, "10 · Calidad y reproducibilidad", "Código verificable, contratos visibles y despliegue repetible.", "FastAPI, React/Vite, workers Python e imágenes ARM64 conectados por infraestructura como código.")
    add_metric_card(slide, "19", "Backend passed", "1 integración Redis skipped", 0.62, 1.78, 1.84, "white", "teal")
    add_metric_card(slide, "9 / 9", "Frontend tests", "Vitest en ejecución única", 2.62, 1.78, 1.84, "white", "purple")
    add_metric_card(slide, "ARM64", "ECS / ECR", "Backend y workers", 0.62, 3.10, 1.84, "white", "aws")
    add_metric_card(slide, "23", "Fases IaC", "CloudFormation 00–22", 2.62, 3.10, 1.84, "white", "blue")
    add_box(slide, 0.62, 4.58, 3.84, 1.52, "navy", "navy", radius=True)
    add_text(slide, "VALIDACIONES", 0.88, 4.86, 2.2, 0.22, 9.2, "mint", True, MONO_FONT)
    add_text(slide, "Vite build · compileall · diagnósticos · diff check · endpoints 200", 0.88, 5.19, 3.25, 0.63, 12.5, "white", True, DISPLAY_FONT)
    add_box(slide, 4.78, 1.74, 7.90, 4.72, "white", "line", radius=True)
    add_picture_cover(slide, CAPTURES / "Screenshot_14.png", 4.90, 1.86, 7.66, 4.48, centering=(0.49, 0.42), pixels_per_inch=220)
    add_chip(slide, "OPENAPI 3.1", 10.80, 1.98, 1.42, "soft_green", "green", "green", 8.7)
    add_footer(slide, 11)
    add_notes(slide, """
Además de funcionar, el sistema es verificable. El backend registra 19 pruebas aprobadas y una integración Redis omitida; el frontend tiene 9 de 9. Las imágenes se construyen para ARM64 y Swagger expone los contratos. Build, compilación, diagnósticos y comprobaciones Git completan la trazabilidad entre código y despliegue.
""")

    # 12 — Cost
    slide = new_slide(prs, "navy")
    add_title(slide, "11 · Coste y gobierno", "Staging optimizado para aprender sin perder control financiero.", "Rangos de planificación en us-east-1; no son una factura ni incluyen Bedrock/RAG.", dark=True)
    cost_cards = [
        ("USD 35–55", "FOUNDATION / MES", "NAT · RDS · Redis · ECR", "blue"),
        ("USD 95–120", "STAGING / MES", "Plataforma ARM64 sin Bedrock", "aws"),
        ("USD 9.37–11.84", "PRUEBA 72 HORAS", "Estimación base de bajo tráfico", "teal"),
    ]
    for idx, (value, label, detail, accent) in enumerate(cost_cards):
        x = 0.62 + idx * 4.08
        add_box(slide, x, 1.86, 3.72, 1.65, "navy_2", "slate", radius=True)
        add_line(slide, x, 1.86, 3.72, accent, 0.065)
        add_text(slide, value, x + 0.25, 2.17, 3.22, 0.47, 23, "white", True, DISPLAY_FONT)
        add_text(slide, label, x + 0.25, 2.75, 3.22, 0.22, 9.2, accent, True, MONO_FONT)
        add_text(slide, detail, x + 0.25, 3.08, 3.22, 0.22, 10.2, "muted_dark")
    add_text(slide, "PALANCAS DE OPTIMIZACIÓN", 0.68, 4.02, 4.0, 0.22, 9.5, "mint", True, MONO_FONT)
    levers = [
        ("ARM64 / Graviton", "Mejor relación precio-rendimiento en Fargate.", "teal"),
        ("NAT instance", "Menor coste de staging; no equivale a alta disponibilidad.", "blue"),
        ("Workers opt-in", "Notificaciones 0/0 y acciones automáticas apagadas.", "purple"),
        ("EC2 reutilizada", "Productor sintético con coste marginal casi cero.", "aws"),
    ]
    for idx, (title, body, accent) in enumerate(levers):
        add_small_card(slide, title, body, 0.62 + idx * 3.08, 4.38, 2.78, 1.46, accent, dark=True)
    add_box(slide, 0.62, 6.15, 12.08, 0.42, "navy_3", "slate", radius=True)
    add_text(slide, "EXCLUIDO DEL RANGO BASE: Bedrock · Knowledge Base · S3 Vectors · tráfico significativo · HA avanzada", 0.82, 6.24, 11.68, 0.24, 9.0, "muted_dark", True, MONO_FONT, PP_ALIGN.CENTER)
    add_footer(slide, 12, dark=True)
    add_notes(slide, """
El objetivo es gobierno, no una cifra artificialmente perfecta. La plataforma completa de staging se estima entre 95 y 120 dólares mensuales antes de Bedrock y vectores. Para 72 horas, la base está alrededor de 9 a 12 dólares. ARM64, NAT instance y workers opt-in reducen gasto, pero la NAT actual no ofrece alta disponibilidad.
""")

    # 13 — Bedrock
    slide = new_slide(prs, "paper")
    add_title(slide, "12 · Bedrock: límite actual", "La IA generativa está preparada, no sobrevendida.", "La cuenta devuelve NOT_AUTHORIZED; el flujo productivo conserva un fallback determinístico seguro.")
    add_box(slide, 0.62, 1.78, 5.00, 1.20, "soft_red", "red", radius=True)
    add_text(slide, "amazon.nova-lite-v1:0", 0.90, 2.02, 3.18, 0.28, 12.5, "ink", True, MONO_FONT)
    add_chip(slide, "NOT_AUTHORIZED", 4.05, 2.02, 1.28, "red", "white", size=8.4)
    add_text(slide, "Proveedor generativo de explicaciones", 0.90, 2.48, 3.65, 0.22, 10.2, "muted")
    add_box(slide, 0.62, 3.16, 5.00, 1.20, "soft_amber", "amber", radius=True)
    add_text(slide, "amazon.titan-embed-text-v2:0", 0.90, 3.40, 3.45, 0.28, 11.5, "ink", True, MONO_FONT)
    add_chip(slide, "NOT_AUTHORIZED", 4.05, 3.40, 1.28, "amber", "white", size=8.4)
    add_text(slide, "Embeddings para Knowledge Base / RAG", 0.90, 3.86, 3.65, 0.22, 10.2, "muted")
    add_box(slide, 0.62, 4.68, 5.00, 1.46, "white", "line", radius=True)
    add_text(slide, "CAUSA PROBABLE, NO CONFIRMADA", 0.90, 4.94, 3.95, 0.23, 9.2, "amber", True, MONO_FONT)
    add_text(slide, "Cuenta nueva, acceso regional de modelos, IAM, primera activación o facturación deben verificarse con Model Access y CloudTrail.", 0.90, 5.30, 4.32, 0.58, 11.2, "muted")
    add_box(slide, 5.96, 1.78, 6.74, 4.36, "navy", "navy", radius=True)
    add_text(slide, "FALLBACK SEGURO", 6.28, 2.06, 2.6, 0.23, 9.4, "mint", True, MONO_FONT)
    fallback_nodes = [
        ("Bedrock no disponible", "red"),
        ("Reglas determinísticas", "purple"),
        ("AIAnalysis + Alert", "teal"),
    ]
    for idx, (label, accent) in enumerate(fallback_nodes):
        x = 6.32 + idx * 2.04
        add_box(slide, x, 2.58, 1.63, 0.78, "navy_3", accent, radius=True)
        add_text(slide, label, x + 0.08, 2.71, 1.47, 0.48, 10.2, "white", True, BODY_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        if idx < 2:
            add_arrow(slide, x + 1.69, 2.87, 0.24, 0.18, "mint")
    add_line(slide, 6.30, 3.76, 5.95, "slate", 0.012)
    add_bullets(
        slide,
        [
            "Lex V2 continúa funcionando de manera independiente.",
            "No hay reintentos infinitos ni afirmaciones de RAG activo.",
            "AI_ENABLE_ACTIONS=false mantiene el sistema en sólo lectura.",
            "Próximo paso: autorización, pruebas, límites de tokens y budgets.",
        ],
        6.30,
        4.05,
        5.94,
        1.58,
        11.3,
        "muted_dark",
        "mint",
        4,
    )
    add_footer(slide, 13)
    add_notes(slide, """
Aquí demostramos rigor técnico. Dos modelos devuelven NOT_AUTHORIZED. La cuenta nueva es una causa probable, pero no podemos declararla como causa raíz sin revisar Model Access, IAM, Billing y CloudTrail. El diseño degrada con seguridad: las reglas siguen creando análisis y alertas, Lex continúa independiente y no se ejecutan acciones.
""")

    # 14 — Roadmap
    slide = new_slide(prs, "navy")
    add_title(slide, "13 · Roadmap", "De staging validado a plataforma de producción.", "La modularidad permite promover capacidades sin reescribir el núcleo.", dark=True)
    roadmap = [
        ("01", "BORDE SEGURO", "S3 privado\nCloudFront + OAC\nRoute 53 + ACM\nALB HTTPS", "blue"),
        ("02", "RESILIENCIA", "WAF y rate limits\nRDS Multi-AZ\nRedis failover\nAutoscaling", "teal"),
        ("03", "GENAI / RAG", "Autorizar Bedrock\nKnowledge Base\nS3 Vectors\nRedacción + budgets", "purple"),
        ("04", "RESPUESTA GOBERNADA", "Canales aprobados\nRetry + DLQ\nAuditoría\nAprobación + rollback", "aws"),
    ]
    for idx, (number, title, body, accent) in enumerate(roadmap):
        x = 0.62 + idx * 3.08
        add_box(slide, x, 1.92, 2.78, 3.54, "navy_2", "slate", radius=True)
        add_icon_circle(slide, number, x + 0.22, 2.18, 0.54, accent, "white" if accent != "aws" else "navy", 10)
        add_text(slide, title, x + 0.22, 2.93, 2.34, 0.48, 13.2, "white", True, DISPLAY_FONT)
        add_text(slide, body, x + 0.22, 3.64, 2.34, 1.38, 12.0, "muted_dark")
        if idx < 3:
            add_arrow(slide, x + 2.83, 3.46, 0.22, 0.18, "mint")
    add_box(slide, 1.82, 5.85, 9.69, 0.61, "navy_3", "mint", radius=True)
    add_rich_text(
        slide,
        [
            ("VISIÓN  →  ", 11, "mint", True, MONO_FONT),
            ("más disponibilidad, más inteligencia y más automatización, ", 12.5, "white", False, BODY_FONT),
            ("siempre con control humano y evidencia.", 12.5, "aws", True, BODY_FONT),
        ],
        2.08,
        6.05,
        9.15,
        0.28,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 14, dark=True)
    add_notes(slide, """
El roadmap empieza por asegurar el borde y eliminar HTTP. Después incorpora alta disponibilidad y autoscaling. Sólo entonces activamos Bedrock y RAG con autorización, redacción y control de costes. Finalmente habilitamos notificaciones y acciones bajo auditoría, aprobación y rollback. El punto clave es que el núcleo asíncrono no necesita reconstruirse.
""")

    # 15 — Close
    slide = new_slide(prs, "navy")
    add_picture_cover(slide, CAPTURES / "Screenshot_17.png", 7.20, 0, 6.133, 7.5, centering=(0.68, 0.44), pixels_per_inch=210)
    add_box(slide, 6.93, 0, 0.28, 7.5, "teal", radius=False)
    add_logo_pill(slide, AWS_LOGO, 0.68, 0.48, 1.32, 0.55)
    add_logo_pill(slide, CODIGO_LOGO, 2.14, 0.48, 2.20, 0.55)
    add_text(slide, "CIERRE", 0.70, 1.46, 1.8, 0.24, 10, "mint", True, MONO_FONT)
    add_text(slide, "Señales que terminan\nen decisiones.", 0.68, 1.86, 5.95, 1.34, 31, "white", True, DISPLAY_FONT)
    add_text(slide, "SentinelMonitorIA demuestra que una plataforma AIOps puede ser útil hoy, honesta sobre sus límites y estar preparada para evolucionar mañana.", 0.70, 3.42, 5.78, 0.93, 14.2, "muted_dark")
    add_chip(slide, "DESPLEGADO", 0.70, 4.60, 1.42, "green", "white")
    add_chip(slide, "TRAZABLE", 2.28, 4.60, 1.30, "teal", "white")
    add_chip(slide, "EXTENSIBLE", 3.74, 4.60, 1.36, "purple", "white")
    demo_button = add_box(slide, 0.70, 5.34, 2.36, 0.64, "aws", "aws", radius=True)
    demo_button.click_action.hyperlink.address = DEMO_URL
    add_text(slide, "ABRIR DEMO AWS  ↗", 0.70, 5.43, 2.36, 0.42, 10.5, "navy", True, MONO_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    repo_button = add_box(slide, 3.24, 5.34, 2.70, 0.64, "navy_3", "mint", radius=True)
    repo_button.click_action.hyperlink.address = REPO_URL
    add_text(slide, "REPOSITORIO / DOCS  ↗", 3.24, 5.43, 2.70, 0.42, 9.7, "mint", True, MONO_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(slide, "Jeffersson Pretell Velasquez · Fernanda Flórez Hereña · Jose Jose Ramones Moreno", 0.70, 6.42, 5.95, 0.38, 9.4, "muted_dark", False, BODY_FONT)
    add_text(slide, "GRACIAS · PREGUNTAS", 0.70, 6.92, 3.20, 0.20, 8.8, "mint", True, MONO_FONT)
    add_text(slide, "15 / 15", 6.08, 6.92, 0.58, 0.20, 8.5, "muted_dark", True, MONO_FONT, PP_ALIGN.RIGHT)
    add_notes(slide, f"""
Cerrar con tres pruebas: está desplegado, es trazable y puede evolucionar. Invitar al jurado a abrir la demo o revisar el repositorio. No mostrar credenciales en la diapositiva; el acceso controlado se entrega por el dossier de evaluación. Demo: {DEMO_URL}\nHealth: {HEALTH_URL}\nRepositorio: {REPO_URL}
""")

    return prs


def validate_presentation(path: Path) -> dict[str, int]:
    """Validate structure, editability, bounds, links, notes and secret hygiene."""
    import tempfile
    import zipfile

    if not path.exists() or path.stat().st_size < 1_000_000:
        raise ValueError("The generated presentation is missing or unexpectedly small")

    with zipfile.ZipFile(path) as archive:
        damaged_member = archive.testzip()
        if damaged_member is not None:
            raise ValueError(f"Damaged Open XML member: {damaged_member}")
        members = archive.namelist()
        media = [name for name in members if name.startswith("ppt/media/")]
        relationship_xml = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in members
            if name.endswith(".rels")
        )
        content_xml = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in members
            if name.endswith(".xml")
        )

    forbidden_values = (
        "S3ntinel!Demo2026",
        "eyJhbGciOi",
        "qoZHpYKu2l3HNfuXCQGTcx2qIlrQ32a07ELBIZ13zZ4",
    )
    exposed = [value for value in forbidden_values if value in content_xml or value in relationship_xml]
    if exposed:
        raise ValueError(f"Sensitive value found in presentation XML: {exposed[0]}")
    if DEMO_URL not in relationship_xml or REPO_URL not in relationship_xml:
        raise ValueError("Expected demo and repository hyperlinks were not embedded")

    presentation = Presentation(path)
    if len(presentation.slides) != TOTAL_SLIDES:
        raise ValueError(f"Expected {TOTAL_SLIDES} slides, found {len(presentation.slides)}")
    aspect_ratio = presentation.slide_width / presentation.slide_height
    if abs(aspect_ratio - (16 / 9)) > 0.002:
        raise ValueError(f"Unexpected aspect ratio: {aspect_ratio:.4f}")

    out_of_bounds: list[str] = []
    missing_notes: list[int] = []
    empty_slides: list[int] = []
    picture_count = 0
    text_shape_count = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if not notes:
            missing_notes.append(slide_number)
        slide_text_shapes = 0
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > presentation.slide_width
                or shape.top + shape.height > presentation.slide_height
            ):
                out_of_bounds.append(f"slide {slide_number}: {shape.name}")
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                picture_count += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shape_count += 1
                slide_text_shapes += 1
        if slide_text_shapes == 0:
            empty_slides.append(slide_number)

    if out_of_bounds:
        raise ValueError("Shapes outside slide bounds: " + ", ".join(out_of_bounds))
    if missing_notes:
        raise ValueError(f"Slides without presenter notes: {missing_notes}")
    if empty_slides:
        raise ValueError(f"Slides without editable text: {empty_slides}")
    if picture_count < 12 or len(media) < 8:
        raise ValueError("Expected visual evidence was not embedded")

    # Save and reopen a temporary copy to catch package/relationship problems.
    with tempfile.TemporaryDirectory(prefix="sentinel-pptx-") as temporary_directory:
        roundtrip = Path(temporary_directory) / "roundtrip.pptx"
        presentation.save(roundtrip)
        reopened = Presentation(roundtrip)
        if len(reopened.slides) != TOTAL_SLIDES:
            raise ValueError("Round-trip validation changed the slide count")

    return {
        "slides": len(presentation.slides),
        "pictures": picture_count,
        "media": len(media),
        "notes": TOTAL_SLIDES - len(missing_notes),
        "text_shapes": text_shape_count,
        "bytes": path.stat().st_size,
    }


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_presentation()
    presentation.save(OUTPUT)
    results = validate_presentation(OUTPUT)
    print(f"Created and validated: {OUTPUT}")
    print(f"Slides: {results['slides']} · notes: {results['notes']}")
    print(f"Pictures: {results['pictures']} · media files: {results['media']}")
    print(f"Editable text shapes: {results['text_shapes']}")
    print(f"Size: {results['bytes']:,} bytes")


if __name__ == "__main__":
    main()
